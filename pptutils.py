from pptx import Presentation
from utils import read_data


def generate(papers,output):
    ppt = Presentation(f'./data/slide.pptx')

    slides = ppt.slides
    for n, paper in enumerate(papers):
        slide = slides[n]
        for i, shape in enumerate(slide.shapes):
            def set_format(target_run, font):
                target_run.font.name = font.name
                target_run.font.size = font.size
                target_run.font.color.theme_color = font.color.theme_color

            if shape.has_text_frame and len(shape.text_frame.paragraphs) > 1:
                for j, paragraph in enumerate(shape.text_frame.paragraphs):
                    if j == 1:
                        paragraph.runs[0].text = f'{paper.title}'
                    elif j == 3:
                        font = paragraph.runs[0].font
                        paragraph.clear()
                        assert len(paper.authors) == len(paper.author_homepages)
                        for index, t in enumerate(paper.authors):
                            run = paragraph.add_run()
                            set_format(run, font)
                            if index == len(paper.authors) - 1:
                                run.text = f'{paper.authors[index]},{paper.author_homepages[index]}'
                            else:
                                run.text = f'{paper.authors[index]},{paper.author_homepages[index]}\n'
                    elif j == 5:
                        font = paragraph.runs[0].font
                        paragraph.clear()
                        run = paragraph.add_run()
                        set_format(run,font)
                        run.text = paper.symbol
                    elif j == 7:
                        font = paragraph.runs[0].font
                        paragraph.clear()
                        run = paragraph.add_run()
                        set_format(run,font)
                        run.text = paper.content
                    elif j == 9:
                        font = paragraph.runs[0].font
                        paragraph.clear()
                        run = paragraph.add_run()
                        set_format(run,font)
                        run.text = paper.year
    ppt.save(output)


if __name__ == '__main__':
    papers = read_data('./data/combined_result.csv')
    generate(papers,'./test.pptx')